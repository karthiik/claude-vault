#!/usr/bin/env python3
"""
Lockton PowerPoint Deck Generator

Generates Lockton-branded presentations using the official template.
Usage:
    python3 generate_deck.py --template TEMPLATE --output OUTPUT --config CONFIG_JSON

Or import and use directly:
    from generate_deck import LocktonDeck
    deck = LocktonDeck()
    deck.add_title_slide("My Presentation", "Subtitle")
    deck.add_content_slide("Topic", ["Point 1", "Point 2"])
    deck.save("output.pptx")
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Run: pip3 install python-pptx")
    sys.exit(1)


# Default template path (relative to vault root)
DEFAULT_TEMPLATE = "1-Projects/2026 PACE/PACE SteerCo Update - December 25.pptx"

# Layout indices in Lockton template
LAYOUTS = {
    "title_dark": 15,       # Title Only_Black - for cover slides
    "title_light": 16,      # Title Only - for section dividers
    "content_dark": 17,     # Title and Content_Black
    "content_light": 18,    # Title and Content_White
    "one_column": 22,       # One-Column Content
    "copyright": 19,        # Copyright - Standard logo
}


class LocktonDeck:
    """Lockton-branded PowerPoint generator"""

    def __init__(self, template_path=None):
        """Initialize with Lockton template"""
        self.template_path = template_path or DEFAULT_TEMPLATE
        self.prs = Presentation(self.template_path)

        # Clear existing slides but keep layouts
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def add_title_slide(self, title, subtitle=None, layout="title_dark"):
        """Add a title/cover slide"""
        layout_idx = LAYOUTS.get(layout, LAYOUTS["title_dark"])
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Find subtitle placeholder if exists
        if subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Subtitle
                    shape.text = subtitle
                    break

        return slide

    def add_section_slide(self, title, subtitle=None):
        """Add a section divider slide"""
        return self.add_title_slide(title, subtitle, layout="title_light")

    def add_content_slide(self, title, bullets=None, body=None, layout="content_light"):
        """Add a content slide with bullets or body text"""
        layout_idx = LAYOUTS.get(layout, LAYOUTS["content_light"])
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Find body placeholder
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in [10, 13]:  # Body placeholders
                body_placeholder = shape
                break

        if body_placeholder:
            if bullets:
                body_placeholder.text = "\n".join(f"• {b}" for b in bullets)
            elif body:
                body_placeholder.text = body

        return slide

    def add_one_column_slide(self, title, content):
        """Add a one-column content slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUTS["one_column"]])

        if slide.shapes.title:
            slide.shapes.title.text = title

        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 13:
                shape.text = content
                break

        return slide

    def add_two_column_slide(self, title, left_content, right_content):
        """Add a two-column comparison slide (using content_light with custom shapes)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUTS["content_light"]])

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Clear default body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in [10, 13]:
                shape.text = ""

        # Add left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True

        # Left header
        p = tf.paragraphs[0]
        p.text = left_content.get("header", "")
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 51, 102)

        # Left items
        for item in left_content.get("items", []):
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)

        # Add right column
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True

        # Right header
        p = tf.paragraphs[0]
        p.text = right_content.get("header", "")
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 163, 173)  # Lockton teal

        # Right items
        for item in right_content.get("items", []):
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)

        return slide

    def add_table_slide(self, title, headers, rows, quote=None):
        """Add a slide with a table"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUTS["content_light"]])

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Clear body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in [10, 13]:
                shape.text = ""

        y_start = Inches(1.5)

        # Add quote if provided
        if quote:
            quote_box = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(12.3), Inches(0.7))
            tf = quote_box.text_frame
            p = tf.paragraphs[0]
            p.text = f'"{quote}"'
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.alignment = PP_ALIGN.CENTER
            y_start = Inches(2.3)

        # Create table
        cols = len(headers)
        num_rows = len(rows) + 1
        table = slide.shapes.add_table(
            num_rows, cols,
            Inches(0.5), y_start,
            Inches(12.3), Inches(0.5 * num_rows)
        ).table

        # Style header row
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Navy
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Data rows
        for r_idx, row in enumerate(rows):
            for c_idx, text in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(text)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r_idx % 2 == 0 else RGBColor(240, 240, 240)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)

        return slide

    def add_closing_slide(self):
        """Add copyright/closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUTS["copyright"]])
        return slide

    def save(self, output_path):
        """Save the presentation"""
        self.prs.save(output_path)
        print(f"✓ Saved: {output_path}")
        return output_path

    @classmethod
    def from_config(cls, config, template_path=None):
        """Create deck from JSON config"""
        deck = cls(template_path)

        for slide_config in config.get("slides", []):
            layout = slide_config.get("layout", 18)
            title = slide_config.get("title", "")

            if layout in [15, LAYOUTS["title_dark"]]:
                deck.add_title_slide(title, slide_config.get("subtitle"))
            elif layout in [16, LAYOUTS["title_light"]]:
                deck.add_section_slide(title, slide_config.get("subtitle"))
            elif layout in [17, 18, LAYOUTS["content_dark"], LAYOUTS["content_light"]]:
                deck.add_content_slide(
                    title,
                    bullets=slide_config.get("bullets"),
                    body=slide_config.get("body"),
                    layout="content_dark" if layout == 17 else "content_light"
                )
            elif layout in [22, LAYOUTS["one_column"]]:
                deck.add_one_column_slide(title, slide_config.get("body", ""))
            elif layout in [19, LAYOUTS["copyright"]]:
                deck.add_closing_slide()
            elif slide_config.get("type") == "table":
                deck.add_table_slide(
                    title,
                    slide_config.get("headers", []),
                    slide_config.get("rows", []),
                    slide_config.get("quote")
                )
            elif slide_config.get("type") == "two_column":
                deck.add_two_column_slide(
                    title,
                    slide_config.get("left", {}),
                    slide_config.get("right", {})
                )

        return deck


def main():
    parser = argparse.ArgumentParser(description="Generate Lockton-branded PowerPoint")
    parser.add_argument("--template", default=DEFAULT_TEMPLATE, help="Template path")
    parser.add_argument("--output", required=True, help="Output file path")
    parser.add_argument("--config", required=True, help="JSON config file path")

    args = parser.parse_args()

    # Load config
    with open(args.config) as f:
        config = json.load(f)

    # Generate deck
    deck = LocktonDeck.from_config(config, args.template)
    deck.save(args.output)


if __name__ == "__main__":
    main()
