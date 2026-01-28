#!/usr/bin/env python3
"""
Innovation Pitch Deck Generator
Creates a professional PowerPoint presentation for Byron Clymer
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Color palette (Professional blues)
DARK_BLUE = RGBColor(0, 51, 102)
ACCENT_BLUE = RGBColor(0, 102, 153)
LIGHT_BLUE = RGBColor(204, 229, 255)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(240, 240, 240)
GREEN = RGBColor(46, 125, 50)
ORANGE = RGBColor(230, 126, 34)
PURPLE = RGBColor(106, 27, 154)
TEAL = RGBColor(0, 121, 107)


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with dark blue background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_header(slide, prs, title):
    """Add consistent header to slides"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_quote_box(slide, quote, y_pos):
    """Add a styled quote box"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), y_pos, Inches(9.2), Inches(0.75))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.fill.background()

    text = slide.shapes.add_textbox(Inches(0.6), y_pos + Inches(0.12), Inches(8.8), Inches(0.5))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, quote=None):
    """Add a content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    y_start = Inches(1.3)
    if quote:
        add_quote_box(slide, quote, y_start)
        y_start = Inches(2.2)

    content = slide.shapes.add_textbox(Inches(0.4), y_start, Inches(9.2), Inches(3.5))
    tf = content.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)

    return slide


def add_table_slide(prs, title, headers, rows, quote=None):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    y_start = Inches(1.3)
    if quote:
        add_quote_box(slide, quote, y_start)
        y_start = Inches(2.2)

    cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(num_rows, cols, Inches(0.4), y_start, Inches(9.2), Inches(0.45 * num_rows)).table

    # Set column widths
    if cols == 2:
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(5.2)
    elif cols == 3:
        table.columns[0].width = Inches(2.8)
        table.columns[1].width = Inches(3.4)
        table.columns[2].width = Inches(3)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, text in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY

    return slide


def add_flow_slide(prs, title, items, caption=None):
    """Add a horizontal flow diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    colors = [GREEN, ORANGE, ACCENT_BLUE]
    num = len(items)
    box_w = Inches(2.5)
    box_h = Inches(1.4)
    gap = Inches(0.5)
    total_w = num * box_w + (num - 1) * gap
    start_x = (prs.slide_width - total_w) / 2
    y = Inches(2.2)

    for i, item in enumerate(items):
        x = start_x + i * (box_w + gap)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Text
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Arrow
        if i < num - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.08),
                y + box_h / 2 - Inches(0.2),
                Inches(0.35),
                Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.4), Inches(4.2), Inches(9.2), Inches(0.6))
        tf = cap.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_progression_slide(prs, title, levels, note=None):
    """Add a progression/levels diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    start_y = Inches(1.4)
    box_h = Inches(0.65)
    gap = Inches(0.12)

    colors = [
        RGBColor(100, 149, 237),  # Light blue
        RGBColor(65, 105, 225),   # Royal blue
        RGBColor(30, 80, 180),    # Medium blue
        RGBColor(0, 51, 102),     # Dark blue
    ]

    for i, level in enumerate(levels):
        y = start_y + i * (box_h + gap)
        width = Inches(4.5 + i * 0.6)

        box = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.4), y, width, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = level
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    if note:
        note_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(4.2), Inches(3))
        tf = note_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(note.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)

    return slide


def add_icon_grid_slide(prs, title, items):
    """Add a grid of icon-style boxes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    colors = [ACCENT_BLUE, GREEN, ORANGE, PURPLE, TEAL, RGBColor(183, 28, 28)]
    cols = 2
    box_w = Inches(4.4)
    box_h = Inches(1)
    h_gap = Inches(0.3)
    v_gap = Inches(0.15)
    start_x = Inches(0.4)
    start_y = Inches(1.35)

    for i, item in enumerate(items):
        row, col = i // cols, i % cols
        x = start_x + col * (box_w + h_gap)
        y = start_y + row * (box_h + v_gap)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Number
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.08), y + Inches(0.08), Inches(0.35), Inches(0.35))
        num.fill.solid()
        num.fill.fore_color.rgb = WHITE
        num.line.fill.background()
        num.text_frame.paragraphs[0].text = str(i + 1)
        num.text_frame.paragraphs[0].font.size = Pt(12)
        num.text_frame.paragraphs[0].font.bold = True
        num.text_frame.paragraphs[0].font.color.rgb = colors[i % len(colors)]
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Text
        text = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(0.1), Inches(3.8), Inches(0.8))
        tf = text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item["title"]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        if item.get("desc"):
            p2 = tf.add_paragraph()
            p2.text = item["desc"]
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(220, 220, 220)

    return slide


def add_two_column_slide(prs, title, left, right):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    # Left header
    lh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(4.4), Inches(0.45))
    lh.fill.solid()
    lh.fill.fore_color.rgb = ACCENT_BLUE
    lh.line.fill.background()
    lh.text_frame.paragraphs[0].text = left["header"]
    lh.text_frame.paragraphs[0].font.size = Pt(14)
    lh.text_frame.paragraphs[0].font.bold = True
    lh.text_frame.paragraphs[0].font.color.rgb = WHITE
    lh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left content
    lc = slide.shapes.add_textbox(Inches(0.3), Inches(1.85), Inches(4.4), Inches(3))
    tf = lc.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left["items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # Right header
    rh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.3), Inches(4.4), Inches(0.45))
    rh.fill.solid()
    rh.fill.fore_color.rgb = GREEN
    rh.line.fill.background()
    rh.text_frame.paragraphs[0].text = right["header"]
    rh.text_frame.paragraphs[0].font.size = Pt(14)
    rh.text_frame.paragraphs[0].font.bold = True
    rh.text_frame.paragraphs[0].font.color.rgb = WHITE
    rh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right content
    rc = slide.shapes.add_textbox(Inches(5.1), Inches(1.85), Inches(4.4), Inches(3))
    tf = rc.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right["items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    return slide


def add_timeline_slide(prs, title, phases):
    """Add a timeline slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, prs, title)

    colors = [ACCENT_BLUE, GREEN, ORANGE]
    phase_w = Inches(3)
    y_line = Inches(2.5)

    # Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y_line, Inches(9), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_BLUE
    line.line.fill.background()

    for i, phase in enumerate(phases):
        x = Inches(0.5) + i * phase_w

        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.25), y_line - Inches(0.15), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()

        # Label
        label = slide.shapes.add_textbox(x, Inches(1.6), phase_w, Inches(0.5))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = phase["label"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER

        # Items
        items = slide.shapes.add_textbox(x, Inches(2.9), phase_w, Inches(2))
        tf = items.text_frame
        tf.word_wrap = True
        for j, item in enumerate(phase["items"]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)

    return slide


def add_closing_slide(prs, title, one_liner, cta_items):
    """Add closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Quote
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = ACCENT_BLUE
    quote_box.line.fill.background()

    quote = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(8.6), Inches(0.7))
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{one_liner}"'
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # CTA
    cta = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(2.5))
    tf = cta.text_frame
    for i, item in enumerate(cta_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓  {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(18)

    return slide


def create_presentation():
    """Create the full deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1: Title
    add_title_slide(prs, "Innovation Sprint", "Turning Hackathon Energy into Sustainable Results")

    # Slide 2: The Moment
    add_flow_slide(
        prs,
        "The Moment We're In",
        ["Hackathon\nEnergy", "???", "Sustained\nRhythm"],
        "We have momentum. Now we need a system — or the energy fades."
    )

    # Slide 3: MS Playbook
    add_table_slide(
        prs,
        "Microsoft Showed Us the Playbook",
        ["Microsoft's Recipe", "What It Means for Lockton"],
        [
            ["Kaizen First, AI Second", "Find waste → then layer AI"],
            ["Business Leader Sponsorship", '"Unless sponsored, we don\'t work on it"'],
            ["Shared Services Council", "Cross-functional idea prioritization"],
            ["Measure the Impact", "Hard metrics, not just stories"],
        ],
        "When you do Kaizen, you'll find a lot of waste. Then you layer in AI."
    )

    # Slide 4: Byron alignment
    add_table_slide(
        prs,
        "Aligned with Your Strategic Priorities",
        ["Your Priority", "How Innovation Sprints Deliver"],
        [
            ["Scale without doing everything differently", "Repeatable framework, not one-offs"],
            ["Remove friction, cost-constrained", "Kaizen finds waste; AI automates pain"],
            ["Data quality as table stakes", "Focus on governance & clean pipelines"],
            ["Hard metrics behind ROI", "Success criteria defined upfront"],
        ],
        "IT's job is to remove friction — but cost-constrained."
    )

    # Slide 5: Workforce enablement
    add_progression_slide(
        prs,
        "Workforce Enablement Unlocks Scale",
        [
            "Level 1: SUMMARIZE — emails, docs, meetings",
            "Level 2: AUTOMATE — action items, routine tasks",
            "Level 3: CREATE — presentations, reports, drafts",
            "Level 4: TRUST — agentic workflows + oversight",
        ],
        "AI is 'jagged' — brilliant at some things, fails at basics.\n\nYou can't trust it blindly.\n\nIntuition comes from personal use.\n\nWhen workforce is enabled, AI scales naturally."
    )

    # Slide 6: Why Shared Services
    add_icon_grid_slide(
        prs,
        "Why Shared Services First",
        [
            {"title": "High Volume Processes", "desc": "Ideal for Kaizen + AI"},
            {"title": "Cross-Functional Visibility", "desc": "Tim & Nick see impact everywhere"},
            {"title": "Measurable Outcomes", "desc": "Cycle time, errors, cost/transaction"},
            {"title": "Internal Customers", "desc": "Lower risk than client-facing"},
            {"title": "Quick Wins Compound", "desc": "Each win builds the case"},
            {"title": "Council Model", "desc": "Microsoft-proven approach"},
        ]
    )

    # Slide 7: Sprint model
    add_two_column_slide(
        prs,
        "The Sprint Model: 4-Week Cadence",
        {
            "header": "What IT Provides",
            "items": [
                "Infrastructure for secure AI deployment",
                "Data governance & integration layer",
                "Measurement & tracing frameworks",
                "Sprint facilitation & technical enablement"
            ]
        },
        {
            "header": "What Business Provides",
            "items": [
                "Process ownership & domain expertise",
                "Sponsorship & prioritization",
                "Change management for teams",
                "Success metrics & adoption ownership"
            ]
        }
    )

    # Slide 8: Q2 wins
    add_table_slide(
        prs,
        "What Tim and Nick Will See by Q2",
        ["Win", "Metric", "Why It Matters"],
        [
            ["3 Completed Sprints", "End-to-end delivery", "Proven capability"],
            ["Efficiency Gains", "Hours saved, error reduction", "Hard ROI to report"],
            ["Enabled Workforce", "Team AI-fluent", "Adoption not deployment"],
            ["Documented Playbook", "Repeatable process", "Scalability proof"],
        ]
    )

    # Slide 9: Scale path
    add_timeline_slide(
        prs,
        "The Scale Path: Earn the Right",
        [
            {"label": "Q1-Q2: Shared Services", "items": ["Partner with Tim/Nick", "Deliver 3-5 wins", "Build the playbook"]},
            {"label": "Q3: Digital", "items": ["Byron champions model", "Expand to client-facing", "Train the trainers"]},
            {"label": "Q4+: Global", "items": ["Regional CIOs adopt", "Local champions own it", "Becomes Lockton standard"]},
        ]
    )

    # Slide 10: Address concerns
    add_table_slide(
        prs,
        "Addressing the Hard Questions",
        ["Your Concern", "Our Answer"],
        [
            ['"Insurance expects fully baked"', "Complete solutions with human-in-the-loop"],
            ['"Need hard metrics for ROI"', "Success criteria defined before sprint starts"],
            ['"Change management is hard"', "HR partner involved; process owner leads"],
            ['"Skittish on global comms"', "No comms until we have wins to show"],
        ],
        "Bad news is better than uncertainty. Share intent."
    )

    # Slide 11: The Ask
    add_closing_slide(
        prs,
        "The Ask",
        "AI in insurance isn't a technology project — it's a workforce transformation.",
        [
            "Your sponsorship as IT champion",
            "Introduction to Tim and Nick",
            "Air cover for experimentation",
            "Quarterly visibility on progress"
        ]
    )

    # Slide 12: Why Now
    add_content_slide(
        prs,
        "Why This, Why Now",
        [
            "Hackathon momentum — people are energized NOW (energy fades fast)",
            "Microsoft playbook — fresh insights from Seattle, validated approach",
            "Lockton growth targets — $4B → $7.1B requires efficiency, not just headcount",
            "Three forces converging — the window is open"
        ],
        "The choice: Let momentum fade, or channel it into a system that scales."
    )

    output = "Innovation_Pitch_Byron.pptx"
    prs.save(output)
    print(f"✓ Presentation saved: {output}")
    return output


if __name__ == "__main__":
    create_presentation()
